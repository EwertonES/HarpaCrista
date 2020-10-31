import sqlite3
from datetime import datetime
import re
import os

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def get_title(cursor, id):
    cursor.execute(f"SELECT title FROM anthems WHERE idanthem = {id}")
    return cursor.fetchone()[0]


def get_verses(cursor, id):
    cursor.execute(f"SELECT verseorder, verse, ismainverse FROM verses WHERE idanthem = {id} ORDER BY verseorder ASC")
    verse_and_order = cursor.fetchall()
    stanzas = []
    chorus = None
    for item in verse_and_order:
        order, stanza, is_chorus = item
        pat = re.compile("(\d\. \n)?(.*)", re.DOTALL)
        mat = re.match(pat, stanza)
        if mat:
            if is_chorus:
                chorus = mat.group(2).rstrip('\n')
            else:
                stanzas.append(mat.group(2).rstrip('\n'))
    return stanzas, chorus


def reorder_stanzas(stanzas, chorus):
    max_lines = 0
    if chorus:
        lines = chorus.count('\n')
        if lines > max_lines:
            max_lines = lines

    reordered_stanzas = []
    for stanza in stanzas:
        reordered_stanzas.append(stanza)
        lines = stanza.count('\n')
        if lines > max_lines:
            max_lines = lines
        if chorus:
            reordered_stanzas.append(chorus)

    max_lines += 1
    return reordered_stanzas, max_lines


def create_pptx(anthem_id, title, stanzas, max_lines, has_chorus):
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    SLD_BLANK = 6
    WIDTH = Cm(25.4)
    HEIGHT = Cm(19.05)

    # Create Presentation object
    prs = Presentation()

    # Choose slide layout
    slide_layout = prs.slide_layouts[SLD_BLANK]

    # Calculate number of slides
    slide_number = 1
    for stanza in stanzas:
        # Add slide
        slide = prs.slides.add_slide(slide_layout)

        # Change slide
        left = top = width = height = Cm(1)

        ## Create blue rectangle
        blue_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Cm(2.5))
        br_fill = blue_rect.fill
        br_fill.solid()
        br_fill.fore_color.rgb = RGBColor(0xAC, 0xCA, 0xff)

        ## Create white circle
        white_circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(0.25), Cm(0.25), Cm(3.5), Cm(2))
        wc_fill = white_circ.fill
        wc_fill.solid()
        wc_fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

        ## Create title
        if len(title) > 30:
            title_box = slide.shapes.add_textbox(Cm(3.5), Cm(0.5), WIDTH - Cm(3), height)
        else:
            title_box = slide.shapes.add_textbox(Cm(0), Cm(0.5), WIDTH, height)
        title_frame = title_box.text_frame

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER

        title_run = title_paragraph.add_run()
        title_run.text = (title).upper()

        title_font = title_run.font
        title_font.name = "Calibri"
        title_font.size = Pt(28)
        title_font.bold = True
        title_font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        ## Create anthem number
        anthem_id_box = slide.shapes.add_textbox(Cm(0.5), Cm(0.5), Cm(3), height)
        anthem_id_frame = anthem_id_box.text_frame

        anthem_id_paragraph = anthem_id_frame.paragraphs[0]
        anthem_id_paragraph.alignment = PP_ALIGN.CENTER

        anthem_id_run = anthem_id_paragraph.add_run()
        anthem_id_run.text = (str(anthem_id))

        anthem_id_font = anthem_id_run.font
        anthem_id_font.name = "Calibri"
        anthem_id_font.size = Pt(28)
        anthem_id_font.bold = True
        anthem_id_font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        ## Create stanza
        if max_lines < 9:
            lyrics_size = Pt(28)
            lyrics_margin_bottom = HEIGHT - Cm(3)
        else:
            lyrics_size = Pt(22)
            lyrics_margin_bottom = HEIGHT - Cm(3)

        lyrics_id_box = slide.shapes.add_textbox(0, Cm(2.5), WIDTH, lyrics_margin_bottom)
        lyrics_id_frame = lyrics_id_box.text_frame
        lyrics_id_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        lyrics_id_paragraph = lyrics_id_frame.paragraphs[0]
        lyrics_id_paragraph.alignment = PP_ALIGN.CENTER
        
        if max_lines < 9:
            lyrics_id_paragraph.line_spacing = Pt(48)

        lyrics_id_run = lyrics_id_paragraph.add_run()
        lyrics_id_run.text = (stanza.upper())

        lyrics_id_font = lyrics_id_run.font
        lyrics_id_font.name = "Calibri"
        
        lyrics_id_font.size = lyrics_size
        lyrics_id_font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        if has_chorus and slide_number % 2 == 0:
            lyrics_id_font.bold = True

        ## Create ad_logo
        logo = slide.shapes.add_picture('logo2.png', WIDTH - Cm(4), HEIGHT - Cm(3))

        ## Create credits
        anthem_id_box = slide.shapes.add_textbox(0, HEIGHT - Cm(1.5), width, height)
        anthem_id_frame = anthem_id_box.text_frame

        anthem_id_paragraph = anthem_id_frame.paragraphs[0]
        anthem_id_paragraph.alignment = PP_ALIGN.LEFT

        anthem_id_run = anthem_id_paragraph.add_run()
        anthem_id_run.text = ("Harpa Cristã - v0.1\newerton@ewerton.com.br")

        anthem_id_font = anthem_id_run.font
        anthem_id_font.name = "Calibri"
        anthem_id_font.size = Pt(14)
        anthem_id_font.bold = True
        anthem_id_font.color.rgb = RGBColor(0xd3, 0xd3, 0xd3)

        slide_number += 1

    # Save presentation
    prs.save(f'./Harpa Cristã/{anthem_id}. {title}.pptx')


# Create output directory if not exists
if not os.path.exists('Harpa Cristã'):
    os.makedirs('Harpa Cristã')


### GET DATA FROM DATABASE ###
connection = sqlite3.connect('harpa.db')
cursor = connection.cursor()


## GET ANTHEMS
for anthem_id in range(1,641):
    title = get_title(cursor, anthem_id)
    stanzas, chorus = get_verses(cursor, anthem_id)
    
    has_chorus = False
    if chorus:
        has_chorus = True

    reordered_stanzas, max_lines = reorder_stanzas(stanzas, chorus)

    create_pptx(anthem_id, title, reordered_stanzas, max_lines, has_chorus)